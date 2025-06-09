import os
import sys
import pandas as pd
import zipfile
import matplotlib
from docx import Document
from pptx import Presentation

sys.path.insert(0, os.getcwd())

from visualization import generate_charts
from report_docx import generate_docx_report
from report_pptx import generate_pptx_report


def test_reports_include_charts(tmp_path, monkeypatch):
    matplotlib.use("Agg")
    # work in tmp directory to avoid polluting repo
    monkeypatch.chdir(tmp_path)

    session_id = "test_session"
    hw_df = pd.DataFrame({"Tier": ["1"], "Status": ["Active"]})
    sw_df = pd.DataFrame({"Tier": ["1"], "Status": ["Active"]})

    session_folder = os.path.join("temp_sessions", session_id)
    charts = generate_charts(hw_df, sw_df, session_folder)

    docx_path = generate_docx_report(session_id, hw_df, sw_df, charts)
    pptx_path = generate_pptx_report(session_id, hw_df, sw_df, charts)

    with zipfile.ZipFile(docx_path) as zf:
        assert any(name.startswith("word/media/") for name in zf.namelist())

    # verify template content preserved
    doc = Document(docx_path)
    assert doc.paragraphs[0].text == "IT Infrastructure Current State Assessment Report"

    with zipfile.ZipFile(pptx_path) as zf:
        assert any(name.startswith("ppt/media/") for name in zf.namelist())

    ppt = Presentation(pptx_path)
    assert ppt.slides[0].shapes.title.text == "Executive Summary"


def test_reports_hardware_only(tmp_path, monkeypatch):
    matplotlib.use("Agg")
    monkeypatch.chdir(tmp_path)

    session_id = "hw_only"
    hw_df = pd.DataFrame({"Tier": ["1"], "Status": ["Active"]})
    sw_df = None

    session_folder = os.path.join("temp_sessions", session_id)
    charts = generate_charts(hw_df, sw_df, session_folder)

    docx_path = generate_docx_report(session_id, hw_df, sw_df, charts)
    pptx_path = generate_pptx_report(session_id, hw_df, sw_df, charts)

    with zipfile.ZipFile(docx_path) as zf:
        assert any(name.startswith("word/media/") for name in zf.namelist())

    with zipfile.ZipFile(pptx_path) as zf:
        assert any(name.startswith("ppt/media/") for name in zf.namelist())


def test_reports_software_only(tmp_path, monkeypatch):
    matplotlib.use("Agg")
    monkeypatch.chdir(tmp_path)

    session_id = "sw_only"
    hw_df = None
    sw_df = pd.DataFrame({"Tier": ["1"], "Status": ["Active"]})

    session_folder = os.path.join("temp_sessions", session_id)
    charts = generate_charts(hw_df, sw_df, session_folder)

    docx_path = generate_docx_report(session_id, hw_df, sw_df, charts)
    pptx_path = generate_pptx_report(session_id, hw_df, sw_df, charts)

    with zipfile.ZipFile(docx_path) as zf:
        assert any(name.startswith("word/media/") for name in zf.namelist())

    with zipfile.ZipFile(pptx_path) as zf:
        assert any(name.startswith("ppt/media/") for name in zf.namelist())


def test_reports_no_data(tmp_path, monkeypatch):
    matplotlib.use("Agg")
    monkeypatch.chdir(tmp_path)

    session_id = "no_data"
    hw_df = None
    sw_df = None

    session_folder = os.path.join("temp_sessions", session_id)
    charts = generate_charts(hw_df, sw_df, session_folder)

    docx_path = generate_docx_report(session_id, hw_df, sw_df, charts)
    pptx_path = generate_pptx_report(session_id, hw_df, sw_df, charts)

    assert os.path.exists(docx_path)
    assert os.path.exists(pptx_path)
