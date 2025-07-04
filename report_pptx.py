
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

TEMPLATES_DIR = os.path.join(os.path.dirname(__file__), "templates")

def generate_pptx_report(session_id, hw_df, sw_df, chart_paths):
    """Generate an executive summary PPTX report.

    Parameters
    ----------
    session_id : str
        Unique identifier for the current assessment session.
    hw_df : pandas.DataFrame
        Hardware dataframe to summarize.
    sw_df : pandas.DataFrame
        Software dataframe to summarize.
    chart_paths : dict
        Dictionary of chart names to file paths returned by
        :func:`visualization.generate_charts`.

    Returns
    -------
    str or None
        Path to the generated PPTX file or ``None`` if generation failed.
    """
    try:
        output_dir = os.path.join("temp_sessions", session_id)
        os.makedirs(output_dir, exist_ok=True)
        output_path = os.path.join(output_dir, "IT_Current_Status_Executive_Report.pptx")

        template_path = os.path.join(
            TEMPLATES_DIR, "IT_Current_Status_Executive_Report_Template.pptx"
        )
        prs = Presentation(template_path) if os.path.exists(template_path) else Presentation()
        # remove template slides to start clean
        while len(prs.slides) > 0:
            rId = prs.slides._sldIdLst[0].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[0]
        title_slide_layout = prs.slide_layouts[0]
        content_layout = prs.slide_layouts[1]
        blank_layout = prs.slide_layouts[6]

        # Title Slide
        slide = prs.slides.add_slide(title_slide_layout)
        slide.shapes.title.text = "Executive Summary"
        slide.placeholders[1].text = f"Session ID: {session_id}"

        # HW Summary
        slide = prs.slides.add_slide(content_layout)
        slide.shapes.title.text = "Hardware Summary"
        if hw_df is not None and not hw_df.empty:
            if 'Tier' in hw_df.columns:
                tier_text = hw_df['Tier'].value_counts().to_string()
            else:
                tier_text = "No tier data."
            hw_summary = f"Total HW Devices: {len(hw_df)}\nTier Distribution:\n{tier_text}"
        else:
            hw_summary = "No hardware data available."
        slide.placeholders[1].text = hw_summary

        # SW Summary
        slide = prs.slides.add_slide(content_layout)
        slide.shapes.title.text = "Software Summary"
        if sw_df is not None and not sw_df.empty:
            if 'Tier' in sw_df.columns:
                tier_text = sw_df['Tier'].value_counts().to_string()
            else:
                tier_text = "No tier data."
            sw_summary = f"Total SW Packages: {len(sw_df)}\nTier Distribution:\n{tier_text}"
        else:
            sw_summary = "No software data available."
        slide.placeholders[1].text = sw_summary

        # Charts
        for path in chart_paths.values():
            if os.path.exists(path):
                slide = prs.slides.add_slide(blank_layout)
                left = Inches(1)
                top = Inches(1)
                height = Inches(5)
                slide.shapes.add_picture(path, left, top, height=height)

        prs.save(output_path)
        return output_path

    except Exception as e:
        print("❌ Error generating PPTX:", str(e))
        return None
